"""Office-to-Office remediation helpers for DOCX, PPTX, and XLSX files."""

from __future__ import annotations

import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class OfficeRemediationResult:
    input_path: Path
    output_path: Path | None = None
    success: bool = False
    changes: list[str] = field(default_factory=list)
    error_message: str = ""


class OfficeRemediator:
    """Apply deterministic, same-format accessibility improvements to Office files."""

    def __init__(self, llm_client: Any | None = None) -> None:
        self._llm = llm_client

    async def remediate(
        self,
        file_path: Path,
        output_path: Path,
        *,
        title: str = "",
        language: str = "en-US",
    ) -> OfficeRemediationResult:
        result = OfficeRemediationResult(input_path=file_path, output_path=output_path)
        try:
            suffix = file_path.suffix.lower()
            output_path.parent.mkdir(parents=True, exist_ok=True)
            if suffix == ".docx":
                changes = await self._remediate_docx(file_path, output_path, title=title, language=language)
            elif suffix == ".pptx":
                changes = await self._remediate_pptx(file_path, output_path, title=title, language=language)
            elif suffix == ".xlsx":
                changes = await self._remediate_xlsx(file_path, output_path, title=title, language=language)
            else:
                raise ValueError(f"Unsupported Office format: {suffix}")
            result.changes.extend(changes)
            result.success = output_path.exists()
        except Exception as exc:
            result.error_message = str(exc)
        return result

    async def _remediate_docx(
        self,
        input_path: Path,
        output_path: Path,
        *,
        title: str,
        language: str,
    ) -> list[str]:
        from docx import Document
        from docx.oxml import OxmlElement

        doc = Document(str(input_path))
        changes: list[str] = []
        resolved_title = title or input_path.stem.replace("_", " ").strip()
        props = doc.core_properties
        if props.title != resolved_title:
            props.title = resolved_title
            changes.append("Set document title metadata")
        if getattr(props, "language", "") != language:
            props.language = language
            changes.append("Set document language metadata")

        if not self._docx_has_heading_structure(doc):
            heading_candidates = self._docx_heading_candidates(doc)
            if not heading_candidates:
                paragraph = doc.paragraphs[0] if doc.paragraphs else doc.add_paragraph()
                if not paragraph.text.strip():
                    paragraph.text = resolved_title
                heading_candidates = [paragraph]
            for candidate_index, paragraph in enumerate(heading_candidates):
                role = "title" if candidate_index == 0 else "heading"
                style_name = self._apply_docx_structure_style(
                    doc,
                    paragraph,
                    role=role,
                    level=0 if candidate_index == 0 else 1,
                )
                changes.append(
                    f"Marked paragraph '{paragraph.text.strip()[:48]}' as {style_name}"
                )

        for table in doc.tables:
            if not table.rows:
                continue
            tr_pr = table.rows[0]._tr.get_or_add_trPr()
            if tr_pr.find(qn("w:tblHeader")) is None:
                tr_pr.append(OxmlElement("w:tblHeader"))
                changes.append("Marked first row as table header")

        image_index = 0
        for inline_shape in doc.inline_shapes:
            image_index += 1
            alt_text = await self._generate_alt_text(
                self._docx_inline_shape_bytes(inline_shape),
                default=f"{resolved_title} image {image_index}",
            )
            doc_pr = inline_shape._inline.docPr
            if not doc_pr.get("descr"):
                doc_pr.set("descr", alt_text)
                changes.append(f"Added alt text to DOCX image {image_index}")
            if not doc_pr.get("title"):
                doc_pr.set("title", alt_text)

        doc.save(str(output_path))
        return changes

    async def _remediate_pptx(
        self,
        input_path: Path,
        output_path: Path,
        *,
        title: str,
        language: str,
    ) -> list[str]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(input_path))
        changes: list[str] = []
        resolved_title = title or input_path.stem.replace("_", " ").strip()
        props = prs.core_properties
        if props.title != resolved_title:
            props.title = resolved_title
            changes.append("Set presentation title metadata")
        if getattr(props, "language", "") != language:
            props.language = language
            changes.append("Set presentation language metadata")

        for slide_index, slide in enumerate(prs.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape is not None and not title_shape.text.strip():
                candidate = next(
                    (
                        shape.text_frame.text.strip()
                        for shape in slide.shapes
                        if getattr(shape, "has_text_frame", False)
                        and shape is not title_shape
                        and shape.text_frame.text.strip()
                    ),
                    "",
                )
                if candidate:
                    title_shape.text = candidate
                    changes.append(f"Filled missing title on slide {slide_index}")
            elif title_shape is None:
                candidate = next(
                    (
                        shape.text_frame.text.strip().splitlines()[0].strip()
                        for shape in slide.shapes
                        if getattr(shape, "has_text_frame", False)
                        and shape.text_frame.text.strip()
                    ),
                    "",
                )
                if candidate:
                    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, 457200)
                    textbox.text_frame.text = candidate
                    changes.append(f"Added explicit title textbox on slide {slide_index}")

            image_index = 0
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                image_index += 1
                alt_text = await self._generate_alt_text(
                    shape.image.blob,
                    default=f"{resolved_title} slide {slide_index} image {image_index}",
                )
                c_nv_pr = shape._element.nvPicPr.cNvPr
                if not c_nv_pr.get("descr"):
                    c_nv_pr.set("descr", alt_text)
                    changes.append(
                        f"Added alt text to PPTX slide {slide_index} image {image_index}"
                    )
                if not c_nv_pr.get("title"):
                    c_nv_pr.set("title", alt_text)

        prs.save(str(output_path))
        return changes

    async def _remediate_xlsx(
        self,
        input_path: Path,
        output_path: Path,
        *,
        title: str,
        language: str,
    ) -> list[str]:
        from openpyxl import load_workbook

        wb = load_workbook(str(input_path))
        changes: list[str] = []
        resolved_title = title or input_path.stem.replace("_", " ").strip()
        props = wb.properties
        if props.title != resolved_title:
            props.title = resolved_title
            changes.append("Set workbook title metadata")
        if getattr(props, "language", "") != language:
            props.language = language
            changes.append("Set workbook language metadata")

        for ws in wb.worksheets:
            if ws.max_row > 1 and ws.freeze_panes is None:
                ws.freeze_panes = "A2"
                changes.append(f"Enabled header freeze panes on sheet '{ws.title}'")
            if ws.max_row > 1 and ws.max_column > 1 and not ws.auto_filter.ref:
                ws.auto_filter.ref = ws.dimensions
                changes.append(f"Enabled auto-filter on sheet '{ws.title}'")
            if ws.max_row > 1 and ws.print_title_rows is None:
                ws.print_title_rows = "1:1"
                changes.append(f"Set repeating header row on sheet '{ws.title}'")

        wb.save(str(output_path))
        return changes

    def _docx_inline_shape_bytes(self, inline_shape: Any) -> bytes | None:
        embed_ids = inline_shape._inline.xpath(".//*[local-name()='blip']/@*[local-name()='embed']")
        if not embed_ids:
            return None
        rel_id = embed_ids[0]
        part = getattr(inline_shape, "part", None)
        if part is None:
            return None
        image_part = part.related_parts.get(rel_id)
        if image_part is None:
            return None
        return image_part.blob

    async def _generate_alt_text(self, image_bytes: bytes | None, *, default: str) -> str:
        if not image_bytes or self._llm is None or not hasattr(self._llm, "vision"):
            return default[:125]
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            tmp.write(image_bytes)
            tmp_path = Path(tmp.name)
        try:
            text = await self._llm.vision(
                image_path=tmp_path,
                prompt=(
                    "Generate concise, descriptive alt text under 125 characters. "
                    "Respond with only the alt text."
                ),
            )
            text = str(text).strip().strip('"').strip("'")
            return (text or default)[:125]
        except Exception:
            return default[:125]
        finally:
            tmp_path.unlink(missing_ok=True)

    def _docx_has_heading_structure(self, doc: Any) -> bool:
        return any(self._docx_paragraph_has_heading_structure(paragraph) for paragraph in doc.paragraphs)

    def _docx_paragraph_has_heading_structure(self, paragraph: Any) -> bool:
        style_name = (getattr(getattr(paragraph, "style", None), "name", "") or "").strip().lower()
        if style_name.startswith(("title", "heading", "accessibility title", "accessibility heading")):
            return True
        return self._docx_outline_level(paragraph) is not None

    def _docx_heading_candidates(self, doc: Any) -> list[Any]:
        paragraphs = [
            paragraph
            for paragraph in doc.paragraphs
            if paragraph.text.strip() and not self._is_divider_text(paragraph.text)
        ]
        if not paragraphs:
            return []

        candidates: list[Any] = []
        title_block: list[Any] = []
        for paragraph in paragraphs:
            if self._is_body_paragraph(paragraph.text):
                break
            if self._is_heading_like_paragraph(paragraph, next_paragraph=None):
                title_block.append(paragraph)
            if len(title_block) >= 3:
                break

        if not title_block:
            title_block = [paragraphs[0]]
        candidates.extend(title_block[:3])

        for index, paragraph in enumerate(paragraphs):
            if paragraph in candidates:
                continue
            next_paragraph = next((item for item in paragraphs[index + 1 :] if item.text.strip()), None)
            if self._is_heading_like_paragraph(paragraph, next_paragraph=next_paragraph):
                candidates.append(paragraph)

        return candidates

    def _is_heading_like_paragraph(self, paragraph: Any, *, next_paragraph: Any | None) -> bool:
        text = paragraph.text.strip()
        if not text or self._looks_like_url(text):
            return False
        style_name = (getattr(getattr(paragraph, "style", None), "name", "") or "").lower()
        if "list" in style_name:
            return False
        words = text.split()
        if len(words) > 14 or len(text) > 100:
            return False
        if text[-1:] in ".!?;":
            return False
        if not any(char.isalpha() for char in text):
            return False
        if next_paragraph is None:
            return True
        next_text = next_paragraph.text.strip()
        if not next_text:
            return True
        if self._is_body_paragraph(next_text):
            return True
        return len(next_text) > len(text) * 2 or len(next_text.split()) >= len(words) + 6

    def _is_body_paragraph(self, text: str) -> bool:
        stripped = text.strip()
        if not stripped:
            return False
        if self._looks_like_url(stripped) or self._is_divider_text(stripped):
            return False
        words = stripped.split()
        return len(words) >= 15 or len(stripped) >= 100 or stripped[-1:] in ".!?"

    def _is_divider_text(self, text: str) -> bool:
        stripped = text.strip()
        return bool(stripped) and all(char in "_-*=" for char in stripped)

    def _looks_like_url(self, text: str) -> bool:
        lowered = text.strip().lower()
        return lowered.startswith(("http://", "https://", "www."))

    def _apply_docx_structure_style(
        self,
        doc: Any,
        paragraph: Any,
        *,
        role: str,
        level: int,
    ) -> str:
        preferred_styles = ("Title", "Heading 1") if role == "title" else ("Heading 1", "Heading")
        applied_style_name = ""
        for style_name in preferred_styles:
            try:
                paragraph.style = doc.styles[style_name]
                applied_style_name = style_name
                break
            except Exception:
                continue
        if not applied_style_name:
            fallback_name = "Accessibility Title" if role == "title" else "Accessibility Heading 1"
            style = self._ensure_docx_paragraph_style(doc, fallback_name)
            paragraph.style = style
            applied_style_name = style.name
        self._set_docx_outline_level(paragraph, level)
        return applied_style_name

    def _ensure_docx_paragraph_style(self, doc: Any, style_name: str) -> Any:
        from docx.enum.style import WD_STYLE_TYPE

        try:
            return doc.styles[style_name]
        except Exception:
            style = doc.styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
            try:
                style.base_style = doc.styles["Normal"]
            except Exception:
                pass
            return style

    def _set_docx_outline_level(self, paragraph: Any, level: int) -> None:
        from docx.oxml import OxmlElement

        p_pr = paragraph._p.get_or_add_pPr()
        outline = p_pr.find(qn("w:outlineLvl"))
        if outline is None:
            outline = OxmlElement("w:outlineLvl")
            p_pr.append(outline)
        outline.set(qn("w:val"), str(max(level, 0)))

    def _docx_outline_level(self, paragraph: Any) -> int | None:
        p_pr = paragraph._p.pPr
        if p_pr is None:
            return None
        outline = p_pr.find(qn("w:outlineLvl"))
        if outline is None:
            return None
        value = outline.get(qn("w:val"))
        if value is None:
            return None
        try:
            return int(value)
        except (TypeError, ValueError):
            return None


def qn(tag: str) -> str:
    """Qualified XML name for OOXML namespaces."""
    prefix, local = tag.split(":")
    namespaces = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    }
    return f"{{{namespaces[prefix]}}}{local}"
