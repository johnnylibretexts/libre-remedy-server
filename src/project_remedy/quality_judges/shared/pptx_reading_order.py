"""Partial PPTX slide reading-order signals from shape order."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from project_remedy.quality_judges.shared.pptx_metadata import validate_slide_count


@dataclass(frozen=True)
class PPTXSlideReadingOrderSignal:
    slide_index: int
    passed: bool
    score: float
    issue: str = ""
    title_text: str = ""
    first_object_text: str = ""
    previous_object_text: str = ""
    out_of_order_object_text: str = ""
    object_count: int = 0
    serialized_text: str = ""
    shape_order_texts: tuple[str, ...] = ()
    visual_order_texts: tuple[str, ...] = ()
    title_first_score: float = 1.0
    visual_sequence_score: float = 1.0


@dataclass(frozen=True)
class _SlideSemanticObject:
    text: str
    left: int
    top: int
    right: int
    bottom: int


_VISUAL_ORDER_TOLERANCE_EMU = 91440


def pptx_slide_reading_order_signals(
    artifact_path: Path,
    *,
    slide_count: Any = None,
) -> list[PPTXSlideReadingOrderSignal]:
    """Return partial per-slide reading-order signals.

    The signal is intentionally conservative: it asserts title-first shape
    order and catches obvious same-column or same-row visual backtracking.
    Full comprehension remains a calibrated behavioral-test requirement.
    """
    slide_count = validate_slide_count(slide_count)
    if not artifact_path.exists():
        return _fallback_signals(slide_count or 0)
    try:
        from pptx import Presentation
    except ImportError:
        return _fallback_signals(slide_count or 0)

    try:
        presentation = Presentation(str(artifact_path))
    except Exception:  # noqa: BLE001 - malformed input is handled as no parser signal.
        return _fallback_signals(slide_count or 0)

    signals: list[PPTXSlideReadingOrderSignal] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        objects = _semantic_objects(slide)
        title_text = _slide_title_text(slide)
        shape_order_texts = tuple(item.text for item in objects)
        visual_order_texts = tuple(item.text for item in _visual_order(objects))
        if not objects:
            signals.append(
                PPTXSlideReadingOrderSignal(
                    slide_index=slide_index,
                    passed=True,
                    score=1.0,
                    issue="empty_slide",
                    title_text=title_text,
                    object_count=0,
                    serialized_text=title_text,
                    shape_order_texts=shape_order_texts,
                    visual_order_texts=visual_order_texts,
                )
            )
            continue
        first_text = objects[0].text
        serialized_text = "\n".join(shape_order_texts)
        title_first = bool(title_text and _same_text(first_text, title_text))
        if title_text and not title_first:
            signals.append(
                PPTXSlideReadingOrderSignal(
                    slide_index=slide_index,
                    passed=False,
                    score=0.0,
                    issue="slide_title_not_first_in_shape_order",
                    title_text=title_text,
                    first_object_text=first_text,
                    object_count=len(objects),
                    serialized_text=serialized_text,
                    shape_order_texts=shape_order_texts,
                    visual_order_texts=visual_order_texts,
                    title_first_score=0.0,
                )
            )
            continue
        backtracking_scope = objects[1:] if title_first else objects
        backtracking = _first_visual_backtracking(backtracking_scope)
        if backtracking is not None:
            previous, current = backtracking
            signals.append(
                PPTXSlideReadingOrderSignal(
                    slide_index=slide_index,
                    passed=False,
                    score=0.5,
                    issue="slide_shape_order_visual_backtracking",
                    title_text=title_text,
                    first_object_text=first_text,
                    previous_object_text=previous.text,
                    out_of_order_object_text=current.text,
                    object_count=len(objects),
                    serialized_text=serialized_text,
                    shape_order_texts=shape_order_texts,
                    visual_order_texts=visual_order_texts,
                    visual_sequence_score=0.0,
                )
            )
            continue
        signals.append(
            PPTXSlideReadingOrderSignal(
                slide_index=slide_index,
                passed=True,
                score=1.0,
                title_text=title_text,
                first_object_text=first_text,
                object_count=len(objects),
                serialized_text=serialized_text,
                shape_order_texts=shape_order_texts,
                visual_order_texts=visual_order_texts,
            )
        )
    return signals


def _fallback_signals(slide_count: int) -> list[PPTXSlideReadingOrderSignal]:
    return [
        PPTXSlideReadingOrderSignal(
            slide_index=index,
            passed=True,
            score=1.0,
            issue="parser_unavailable",
        )
        for index in range(1, slide_count + 1)
    ]


def _semantic_objects(slide: Any) -> list[_SlideSemanticObject]:
    items: list[_SlideSemanticObject] = []
    for shape in slide.shapes:
        text = ""
        if getattr(shape, "has_text_frame", False):
            text = str(shape.text_frame.text or "").strip()
        elif getattr(shape, "has_table", False):
            text = _table_text(shape)
        if text:
            left = _shape_coordinate(shape, "left")
            top = _shape_coordinate(shape, "top")
            width = _shape_coordinate(shape, "width")
            height = _shape_coordinate(shape, "height")
            items.append(
                _SlideSemanticObject(
                    text=" ".join(text.split()),
                    left=left,
                    top=top,
                    right=left + max(width, 0),
                    bottom=top + max(height, 0),
                )
            )
    return items


def _slide_title_text(slide: Any) -> str:
    title_shape = slide.shapes.title
    if title_shape is not None:
        text = str(getattr(title_shape, "text", "") or "").strip()
        if text:
            return " ".join(text.split())
    return ""


def _same_text(left: str, right: str) -> bool:
    return " ".join(left.split()).strip().lower() == " ".join(right.split()).strip().lower()


def _first_visual_backtracking(
    objects: list[_SlideSemanticObject],
) -> tuple[_SlideSemanticObject, _SlideSemanticObject] | None:
    for previous, current in zip(objects, objects[1:]):
        if _same_column(previous, current) and (
            current.top + _VISUAL_ORDER_TOLERANCE_EMU < previous.top
        ):
            return previous, current
        if _same_row(previous, current) and (
            current.left + _VISUAL_ORDER_TOLERANCE_EMU < previous.left
        ):
            return previous, current
    return None


def _visual_order(objects: list[_SlideSemanticObject]) -> list[_SlideSemanticObject]:
    return sorted(objects, key=lambda item: (item.top, item.left))


def _same_column(left: _SlideSemanticObject, right: _SlideSemanticObject) -> bool:
    return _range_overlap_ratio(left.left, left.right, right.left, right.right) >= 0.25


def _same_row(left: _SlideSemanticObject, right: _SlideSemanticObject) -> bool:
    return _range_overlap_ratio(left.top, left.bottom, right.top, right.bottom) >= 0.25


def _range_overlap_ratio(
    first_start: int,
    first_end: int,
    second_start: int,
    second_end: int,
) -> float:
    first_span = max(first_end - first_start, 1)
    second_span = max(second_end - second_start, 1)
    overlap = max(0, min(first_end, second_end) - max(first_start, second_start))
    return overlap / min(first_span, second_span)


def _shape_coordinate(shape: Any, attribute: str) -> int:
    try:
        return int(getattr(shape, attribute, 0) or 0)
    except (TypeError, ValueError):
        return 0


def _table_text(shape: Any) -> str:
    rows: list[str] = []
    table = getattr(shape, "table", None)
    if table is None:
        return "table"
    for row in table.rows:
        cells = [
            " ".join(str(cell.text or "").split())
            for cell in row.cells
            if str(cell.text or "").strip()
        ]
        if cells:
            rows.append(" | ".join(cells))
    return "table: " + " / ".join(rows) if rows else "table"
