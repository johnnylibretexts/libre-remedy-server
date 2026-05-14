"""PPTX slide-title quality signals."""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from project_remedy.quality_judges.shared.pptx_metadata import validate_slide_count


_GENERIC_TITLES = {
    "agenda",
    "contents",
    "introduction",
    "overview",
    "summary",
    "title",
    "untitled",
}


@dataclass(frozen=True)
class PPTXSlideTitleSignal:
    slide_index: int
    title_text: str
    has_title_placeholder: bool
    issue: str

    @property
    def passed(self) -> bool:
        return self.issue == ""

    @property
    def score(self) -> float:
        if self.issue == "":
            return 1.0
        if self.issue == "duplicate_slide_title":
            return 0.6
        return 0.0


def pptx_slide_title_signals(
    artifact_path: Path,
    *,
    slide_count: Any = None,
) -> list[PPTXSlideTitleSignal]:
    """Return per-slide title quality signals for a PPTX artifact."""
    slide_count = validate_slide_count(slide_count)
    if not artifact_path.exists():
        return _fallback_signals(slide_count)

    try:
        from pptx import Presentation
    except ImportError:
        return _fallback_signals(slide_count)

    try:
        presentation = Presentation(str(artifact_path))
    except Exception:  # noqa: BLE001 - malformed input makes this signal unavailable.
        return _fallback_signals(slide_count)

    provisional: list[tuple[int, str, bool, str]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_text, has_title_placeholder = _slide_title(slide)
        issue = _title_issue(title_text, has_title_placeholder)
        provisional.append((slide_index, title_text, has_title_placeholder, issue))

    title_counts = Counter(
        _normalized_title(title_text)
        for _slide_index, title_text, _has_title_placeholder, issue in provisional
        if title_text and issue == ""
    )
    signals: list[PPTXSlideTitleSignal] = []
    for slide_index, title_text, has_title_placeholder, issue in provisional:
        if issue == "" and title_counts[_normalized_title(title_text)] > 1:
            issue = "duplicate_slide_title"
        signals.append(
            PPTXSlideTitleSignal(
                slide_index=slide_index,
                title_text=title_text,
                has_title_placeholder=has_title_placeholder,
                issue=issue,
            )
        )
    return signals


def _fallback_signals(slide_count: int | None) -> list[PPTXSlideTitleSignal]:
    return [
        PPTXSlideTitleSignal(
            slide_index=index,
            title_text="",
            has_title_placeholder=False,
            issue="parser_unavailable",
        )
        for index in range(1, (slide_count or 0) + 1)
    ]


def _slide_title(slide: Any) -> tuple[str, bool]:
    title_shape = slide.shapes.title
    if title_shape is not None:
        return " ".join((title_shape.text or "").split()), True
    return "", False


def _title_issue(title_text: str, has_title_placeholder: bool) -> str:
    normalized = _normalized_title(title_text)
    if not has_title_placeholder:
        return "missing_slide_title_placeholder"
    if not normalized:
        return "empty_slide_title"
    if normalized in _GENERIC_TITLES or normalized.startswith("slide "):
        return "non_descriptive_slide_title"
    if len(normalized) < 4:
        return "non_descriptive_slide_title"
    return ""


def _normalized_title(value: str) -> str:
    return " ".join(value.split()).strip().lower()
