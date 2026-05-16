from __future__ import annotations

from pathlib import Path
from types import SimpleNamespace
from zipfile import ZipFile

import pytest

from project_remedy.models import FileType
from project_remedy.office_acceptance import (
    OfficeCheckReport,
    OfficeCheckResult,
    OfficeScreenReaderResult,
)
from project_remedy.quality_judges.office import audit as office_audit
from project_remedy.quality_judges.office.docx import (
    DOCXAltTextQualityJudge,
    DOCXComplexContentJudge,
    DOCXDecorativeJudge,
    DOCXHeadingSemanticsJudge,
    DOCXLinkTextJudge,
    DOCXReadingOrderJudge,
    DOCXTableStructureJudge,
)
from project_remedy.quality_judges.office.pptx import (
    PPTXAltTextQualityJudge,
    PPTXComplexContentJudge,
    PPTXDecorativeJudge,
    PPTXHeadingSemanticsJudge,
    PPTXLinkTextJudge,
    PPTXSlideReadingOrderJudge,
    PPTXSlideTitleJudge,
    PPTXTableStructureJudge,
)
from project_remedy.quality_judges.office.xlsx import (
    XLSXAltTextQualityJudge,
    XLSXComplexContentJudge,
    XLSXLinkTextJudge,
    XLSXSheetOrganizationJudge,
    XLSXTableStructureJudge,
)
from project_remedy.quality_judges.shared.base import (
    ModelSeparationError,
    QualityJudgeConfig,
)
from project_remedy.quality_judges.shared.dimensions import DIMENSIONS_BY_FORMAT


def _config() -> QualityJudgeConfig:
    return QualityJudgeConfig(
        backend="ollama",
        model="llama3.1:8b",
        production_models=("gemma4:31b-cloud",),
    )


def _pipeline_config() -> SimpleNamespace:
    return SimpleNamespace(
        api=SimpleNamespace(
            text_model="gemma4:31b-cloud",
            vision_model="gemma4-vision:27b",
            escalation_model="",
            quality_judge_backend="ollama",
            quality_judge_model="llama3.1:8b",
            quality_judge_base_url="",
            behavioral_test_backend="ollama",
            behavioral_test_model="qwen2.5:7b",
            behavioral_test_cache_path="",
        )
    )


def _report(file_type: FileType, results: list[OfficeCheckResult]) -> OfficeCheckReport:
    return OfficeCheckReport(
        file_path=Path("fixture"),
        file_type=file_type,
        results=results,
    )


def _zip_fixture(tmp_path: Path, filename: str, parts: dict[str, str]) -> Path:
    path = tmp_path / filename
    with ZipFile(path, "w") as package:
        for part_name, content in parts.items():
            package.writestr(part_name, content)
    return path


def _mark_docx_header_row(table, value: str = "true") -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    tr_pr = table.rows[0]._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), value)
    tr_pr.append(tbl_header)


def test_docx_judges_return_applicable_dimension_scores() -> None:
    report = _report(
        FileType.DOCX,
        [
            OfficeCheckResult("docx-alt-text", "Images contain alternate text", "Failed"),
            OfficeCheckResult("docx-headings", "Headings exist", "Passed"),
            OfficeCheckResult("docx-table-headers", "Tables have headers", "Passed"),
        ],
    )
    judges = [
        DOCXAltTextQualityJudge(_config()),
        DOCXReadingOrderJudge(_config()),
        DOCXHeadingSemanticsJudge(_config()),
        DOCXTableStructureJudge(_config()),
        DOCXLinkTextJudge(_config()),
        DOCXDecorativeJudge(_config()),
        DOCXComplexContentJudge(_config()),
    ]

    scores = [judge.judge(Path("fixture.docx"), checker_report=report) for judge in judges]

    assert tuple(score.dimension for score in scores) == DIMENSIONS_BY_FORMAT["docx"]
    assert all(score.format == "docx" for score in scores)
    assert scores[0].score == 0.0
    assert scores[0].sample_findings[0]["rule_id"] == "docx-alt-text"
    assert scores[1].dimension == "reading_order"
    assert scores[1].confidence == 0.25
    assert scores[-1].confidence == 0.30


def test_docx_alt_text_judge_reads_ooxml_drawing_metadata(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
                  <wp:docPr id="1" name="Picture 1" descr="Customer journey map" />
                  <wp:docPr id="2" name="Picture 2" />
                </w:document>
            """,
        },
    )

    score = DOCXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.62
    assert score.sample_findings[0]["issue"] == "office_object_missing_alt_text"
    assert score.sample_findings[0]["name"] == "Picture 2"


def test_docx_alt_text_judge_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-generic.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
                  <wp:docPr id="1" name="Picture 1" descr="Image" />
                  <wp:docPr id="2" name="Picture 2"
                    descr="Customer journey map shows 42% renewal path" />
                  <wp:docPr id="3" name="Picture 3"
                    descr="Customer journey map shows 42% renewal path" />
                </w:document>
            """,
        },
    )

    score = DOCXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.3333
    assert score.per_criterion["ooxml_alt_text_presence"] == 1.0
    assert score.per_criterion["ooxml_alt_text_specificity"] == 0.3333
    assert [finding["issue"] for finding in score.sample_findings] == [
        "office_object_non_substitutive_alt_text",
        "office_object_duplicated_substitutive_alt_text",
    ]


def test_docx_heading_semantics_judge_reads_word_heading_styles(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Overview", level=1)
    document.add_heading("Skipped Details", level=3)
    fixture = tmp_path / "headings.docx"
    document.save(fixture)

    score = DOCXHeadingSemanticsJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.70
    assert score.sample_findings[0]["issue"] == "docx_heading_level_skip"
    assert score.sample_findings[0]["level"] == 3


def test_docx_heading_semantics_judge_flags_visual_heading_without_semantic_style(
    tmp_path: Path,
) -> None:
    docx = pytest.importorskip("docx")
    from docx.shared import Pt

    document = docx.Document()
    paragraph = document.add_paragraph()
    run = paragraph.add_run("Financial Overview")
    run.bold = True
    run.font.size = Pt(18)
    document.add_paragraph("The financial trend improved over five years.")
    fixture = tmp_path / "visual-heading.docx"
    document.save(fixture)

    score = DOCXHeadingSemanticsJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.per_criterion["visual_heading_semantics"] == 0.0
    assert score.sample_findings[0]["issue"] == "docx_visual_heading_without_semantic_style"
    assert score.sample_findings[0]["reasons"] == ["large_text", "bold_text"]


def test_docx_table_structure_judge_reads_table_headers(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "$1M"
    _mark_docx_header_row(table)
    fixture = tmp_path / "table.docx"
    document.save(fixture)

    score = DOCXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.confidence == 0.70
    assert score.per_criterion["repeated_header_rows"] == 1.0
    assert score.per_criterion["non_empty_header_cells"] == 0.0
    assert score.sample_findings[0]["issue"] == "docx_table_empty_header_cells"
    assert score.sample_findings[0]["empty_header_columns"] == [1]


def test_docx_table_structure_judge_rejects_false_repeated_header_flag(
    tmp_path: Path,
) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "Region"
    table.cell(1, 0).text = "West"
    _mark_docx_header_row(table, value="false")
    fixture = tmp_path / "false-header.docx"
    document.save(fixture)

    score = DOCXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.per_criterion["repeated_header_rows"] == 0.0
    assert score.sample_findings[0]["issue"] == "docx_table_missing_repeated_header_row"
    assert score.sample_findings[0]["has_header_row"] is False


def test_pptx_judges_cover_slide_title_and_per_slide_reading_order_dimensions() -> None:
    report = _report(
        FileType.PPTX,
        [
            OfficeCheckResult("pptx-alt-text", "Pictures contain alternate text", "Passed"),
            OfficeCheckResult("pptx-slide-titles", "Slides expose titles", "Failed"),
        ],
    )
    judges = [
        PPTXAltTextQualityJudge(_config()),
        PPTXSlideReadingOrderJudge(_config()),
        PPTXHeadingSemanticsJudge(_config()),
        PPTXTableStructureJudge(_config()),
        PPTXLinkTextJudge(_config()),
        PPTXDecorativeJudge(_config()),
        PPTXComplexContentJudge(_config()),
        PPTXSlideTitleJudge(_config()),
    ]

    scores = [
        judge.judge(Path("fixture.pptx"), checker_report=report, slide_count=2)
        for judge in judges
    ]

    assert tuple(score.dimension for score in scores) == DIMENSIONS_BY_FORMAT["pptx"]
    assert scores[2].score == 0.0
    assert scores[7].dimension == "slide_title"
    assert scores[7].sample_findings[0]["rule_id"] == "pptx-slide-titles"
    assert scores[1].confidence == 0.20
    assert scores[1].sample_findings[0]["per_slide"] == [
        {
            "slide_index": 1,
            "score": 1.0,
            "confidence": 0.20,
            "parser_support": "python_pptx_shape_order",
            "passed": True,
            "issue": "parser_unavailable",
            "title_text": "",
            "first_object_text": "",
            "previous_object_text": "",
            "out_of_order_object_text": "",
            "object_count": 0,
            "shape_order_texts": [],
            "visual_order_texts": [],
        },
        {
            "slide_index": 2,
            "score": 1.0,
            "confidence": 0.20,
            "parser_support": "python_pptx_shape_order",
            "passed": True,
            "issue": "parser_unavailable",
            "title_text": "",
            "first_object_text": "",
            "previous_object_text": "",
            "out_of_order_object_text": "",
            "object_count": 0,
            "shape_order_texts": [],
            "visual_order_texts": [],
        },
    ]


def test_pptx_per_slide_judges_reject_invalid_slide_count_metadata() -> None:
    judges = [
        PPTXSlideReadingOrderJudge(_config()),
        PPTXHeadingSemanticsJudge(_config()),
        PPTXSlideTitleJudge(_config()),
    ]

    for judge in judges:
        for value in (True, "2", 1.5, -1):
            with pytest.raises(
                ValueError,
                match="slide_count must be a non-negative integer",
            ):
                judge.judge(Path("fixture.pptx"), slide_count=value)


def test_pptx_reading_order_judge_flags_title_after_body_shape(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Quarterly Results"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "Body first"
    title_element = title._element
    title_element.getparent().remove(title_element)
    slide.shapes._spTree.append(title_element)
    fixture = tmp_path / "bad-order.pptx"
    presentation.save(fixture)

    score = PPTXSlideReadingOrderJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.confidence == 0.45
    assert score.sample_findings[0]["issue"] == "slide_title_not_first_in_shape_order"
    assert score.sample_findings[0]["first_object_text"] == "Body first"


def test_pptx_reading_order_judge_flags_visual_backtracking_in_shape_order(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Results"
    lower = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.5))
    lower.text_frame.text = "Q4 margin details"
    upper = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
    upper.text_frame.text = "Q3 revenue details"
    fixture = tmp_path / "visual-backtracking.pptx"
    presentation.save(fixture)

    score = PPTXSlideReadingOrderJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.per_criterion["shape_order_title_first"] == 1.0
    assert score.per_criterion["shape_order_visual_sequence"] == 0.0
    assert score.sample_findings[0]["issue"] == "slide_shape_order_visual_backtracking"
    assert score.sample_findings[0]["previous_object_text"] == "Q4 margin details"
    assert score.sample_findings[0]["out_of_order_object_text"] == "Q3 revenue details"
    per_slide = score.sample_findings[1]["per_slide"][0]
    assert per_slide["shape_order_texts"] == [
        "Quarterly Results",
        "Q4 margin details",
        "Q3 revenue details",
    ]
    assert per_slide["visual_order_texts"] == [
        "Quarterly Results",
        "Q3 revenue details",
        "Q4 margin details",
    ]


def test_pptx_slide_title_judge_scores_duplicate_and_generic_titles(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")

    presentation = pptx.Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Quarterly Results"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Quarterly Results"
    third = presentation.slides.add_slide(presentation.slide_layouts[5])
    third.shapes.title.text = "Overview"
    fixture = tmp_path / "slide-titles.pptx"
    presentation.save(fixture)

    score = PPTXSlideTitleJudge(_config()).judge(fixture)

    assert score.score == 0.4
    assert score.confidence == 0.65
    assert score.per_criterion["slide_title_presence"] == 1.0
    assert score.per_criterion["slide_title_descriptiveness"] == pytest.approx(2 / 3, abs=0.0001)
    assert score.per_criterion["slide_title_uniqueness"] == pytest.approx(1 / 3, abs=0.0001)
    assert [finding["issue"] for finding in score.sample_findings] == [
        "duplicate_slide_title",
        "duplicate_slide_title",
        "non_descriptive_slide_title",
    ]


def test_pptx_heading_semantics_judge_reads_slide_title_placeholders(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")

    presentation = pptx.Presentation()
    titled = presentation.slides.add_slide(presentation.slide_layouts[5])
    titled.shapes.title.text = "Quarterly Results"
    presentation.slides.add_slide(presentation.slide_layouts[6])
    fixture = tmp_path / "heading-semantics.pptx"
    presentation.save(fixture)

    score = PPTXHeadingSemanticsJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.65
    assert score.per_criterion["slide_title_presence"] == 0.5
    assert score.sample_findings[0]["issue"] == "missing_slide_title_placeholder"
    assert score.sample_findings[0]["slide_index"] == 2


def test_pptx_alt_text_judge_reads_ooxml_drawing_metadata(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                  <p:cSld><p:spTree>
                    <p:pic><p:nvPicPr><p:cNvPr id="1" name="Picture 1"
                      descr="Market share by segment" /></p:nvPicPr></p:pic>
                    <p:pic><p:nvPicPr><p:cNvPr id="2" name="Picture 2" /></p:nvPicPr></p:pic>
                  </p:spTree></p:cSld>
                </p:sld>
            """,
        },
    )

    score = PPTXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.62
    assert score.sample_findings[0]["name"] == "Picture 2"


def test_pptx_alt_text_judge_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-generic.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                  <p:cSld><p:spTree>
                    <p:pic><p:nvPicPr><p:cNvPr id="1" name="Chart 1"
                      descr="Chart" /></p:nvPicPr></p:pic>
                    <p:pic><p:nvPicPr><p:cNvPr id="2" name="Chart 2"
                      descr="Market share rose from 18% to 24%" /></p:nvPicPr></p:pic>
                    <p:pic><p:nvPicPr><p:cNvPr id="3" name="Chart 3"
                      descr="Market share rose from 18% to 24%" /></p:nvPicPr></p:pic>
                  </p:spTree></p:cSld>
                </p:sld>
            """,
        },
    )

    score = PPTXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.3333
    assert score.per_criterion["ooxml_alt_text_presence"] == 1.0
    assert score.per_criterion["ooxml_alt_text_specificity"] == 0.3333
    assert [finding["issue"] for finding in score.sample_findings] == [
        "office_object_non_substitutive_alt_text",
        "office_object_duplicated_substitutive_alt_text",
    ]


def test_xlsx_judges_skip_inapplicable_reading_order_dimension() -> None:
    report = _report(
        FileType.XLSX,
        [
            OfficeCheckResult("xlsx-header-behaviors", "Header navigation aids", "Failed"),
        ],
    )
    judges = [
        XLSXAltTextQualityJudge(_config()),
        XLSXTableStructureJudge(_config()),
        XLSXLinkTextJudge(_config()),
        XLSXComplexContentJudge(_config()),
        XLSXSheetOrganizationJudge(_config()),
    ]

    scores = [
        judge.judge(
            Path("fixture.xlsx"),
            checker_report=report,
            sheet_names=["Sheet1", "Q3 Revenue"],
        )
        for judge in judges
    ]

    assert tuple(score.dimension for score in scores) == DIMENSIONS_BY_FORMAT["xlsx"]
    assert "reading_order" not in {score.dimension for score in scores}
    assert scores[1].score == 0.0
    assert scores[4].score == 0.5
    assert scores[4].sample_findings[0]["sheet_name"] == "Sheet1"


def test_xlsx_sheet_organization_judge_rejects_invalid_sheet_name_metadata() -> None:
    for value in ("Sheet1", [1], [""], True):
        with pytest.raises(
            ValueError,
            match="sheet_names must be a list of non-empty strings",
        ):
            XLSXSheetOrganizationJudge(_config()).judge(
                Path("fixture.xlsx"),
                sheet_names=value,
            )


def test_xlsx_alt_text_judge_reads_drawing_objects(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "drawings.xlsx",
        {
            "xl/drawings/drawing1.xml": """
                <xdr:wsDr
                    xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="1" name="Chart image" descr="Revenue by quarter" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="2" name="Unlabeled image" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                </xdr:wsDr>
            """,
        },
    )

    score = XLSXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.60
    assert score.sample_findings[0]["issue"] == "xlsx_drawing_missing_alt_text"
    assert score.sample_findings[0]["name"] == "Unlabeled image"


def test_xlsx_alt_text_judge_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "drawings-generic.xlsx",
        {
            "xl/drawings/drawing1.xml": """
                <xdr:wsDr
                    xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="1" name="Forecast graph" descr="Graph" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="2" name="Margin chart"
                        descr="Margin improved from 12% to 17%" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="3" name="Duplicate margin chart"
                        descr="Margin improved from 12% to 17%" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                </xdr:wsDr>
            """,
        },
    )

    score = XLSXAltTextQualityJudge(_config()).judge(fixture)

    assert score.score == 0.3333
    assert score.per_criterion["drawing_alt_text_presence"] == 1.0
    assert score.per_criterion["drawing_alt_text_specificity"] == pytest.approx(1 / 3)
    assert [finding["issue"] for finding in score.sample_findings] == [
        "xlsx_drawing_non_substitutive_alt_text",
        "xlsx_drawing_duplicated_substitutive_alt_text",
    ]


def test_xlsx_sheet_organization_judge_uses_sheet_content_terms(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    archive = workbook.create_sheet("Archive")
    archive.append(["Customer", "Renewal"])
    archive.append(["Northwind", "Yes"])
    fixture = tmp_path / "sheets.xlsx"
    workbook.save(fixture)

    score = XLSXSheetOrganizationJudge(_config()).judge(fixture)

    assert score.score == 0.75
    assert score.per_criterion["sheet_name_descriptiveness"] == 1.0
    assert score.per_criterion["sheet_purpose_alignment"] == 0.5
    assert score.sample_findings[0]["issue"] == "sheet_name_purpose_unclear"
    assert score.sample_findings[0]["sheet_name"] == "Archive"


def test_xlsx_sheet_organization_judge_flags_overview_sheet_after_details(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    summary = workbook.create_sheet("Summary")
    summary.append(["Summary", "Revenue"])
    summary.append(["Total", 120000])
    fixture = tmp_path / "summary-after-details.xlsx"
    workbook.save(fixture)

    score = XLSXSheetOrganizationJudge(_config()).judge(fixture)

    assert score.score == 0.75
    assert score.per_criterion["sheet_purpose_alignment"] == 0.5
    assert score.per_criterion["sheet_ordering"] == 0.5
    assert score.sample_findings[0]["issue"] == "overview_sheet_not_first"
    assert score.sample_findings[0]["sheet_name"] == "Summary"


def test_xlsx_sheet_organization_judge_preserves_multiple_sheet_issues(
    tmp_path: Path,
) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    hidden = workbook.create_sheet("Sheet3")
    hidden.append(["Risk", "Score"])
    hidden.append(["Security", 7])
    hidden.sheet_state = "hidden"
    fixture = tmp_path / "hidden-default-sheet.xlsx"
    workbook.save(fixture)

    score = XLSXSheetOrganizationJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.per_criterion["sheet_name_descriptiveness"] == 0.5
    assert score.per_criterion["visible_data_sheets"] == 0.5
    assert [finding["issue"] for finding in score.sample_findings] == [
        "non_descriptive_sheet_name",
        "data_sheet_hidden",
    ]


def test_xlsx_table_structure_judge_reads_excel_table_structures(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.worksheet.table import Table

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["Region", "Revenue"])
    worksheet.append(["West", 120000])
    table = Table(displayName="RevenueTable", ref="A1:B2")
    worksheet.add_table(table)
    fixture = tmp_path / "table-structure.xlsx"
    workbook.save(fixture)

    score = XLSXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.85
    assert score.confidence == 0.70
    assert score.per_criterion["excel_table_presence"] == 1.0
    assert score.per_criterion["header_row_presence"] == 1.0
    assert score.per_criterion["banded_rows"] == 0.0
    assert score.sample_findings[0]["issue"] == "xlsx_table_missing_banded_rows"
    assert score.sample_findings[0]["severity"] == "warning"


def test_xlsx_table_structure_judge_flags_plain_data_ranges(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["Region", "Revenue"])
    worksheet.append(["West", 120000])
    fixture = tmp_path / "plain-range.xlsx"
    workbook.save(fixture)

    score = XLSXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.per_criterion["excel_table_presence"] == 0.0
    assert score.sample_findings[0]["issue"] == "xlsx_data_range_missing_excel_table"
    assert score.sample_findings[0]["sheet_name"] == "Q3 Revenue"


def test_docx_complex_content_judge_scores_chart_descriptions(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "complex.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
                  <wp:docPr id="1" name="Revenue chart" descr="Revenue rose 18% in Q4" />
                  <wp:docPr id="2" name="Cost chart" descr="Cost chart" />
                </w:document>
            """,
        },
    )

    score = DOCXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.55
    assert score.sample_findings[0]["issue"] == "thin_complex_content_description"
    assert score.sample_findings[0]["name"] == "Cost chart"


def test_docx_complex_content_judge_scores_equation_context(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "equations.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">
                  <w:body>
                    <w:p>
                      <w:r><w:t>Conversion equation for quarterly revenue</w:t></w:r>
                      <m:oMath><m:r><m:t>r = q / t</m:t></m:r></m:oMath>
                    </w:p>
                    <w:p>
                      <m:oMath><m:r><m:t>x = y</m:t></m:r></m:oMath>
                    </w:p>
                  </w:body>
                </w:document>
            """,
        },
    )

    score = DOCXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.per_criterion["equation_context"] == 0.5
    assert score.sample_findings[0]["issue"] == "thin_complex_content_description"
    assert score.sample_findings[0]["kind"] == "equation"
    assert score.sample_findings[0]["name"] == "Equation 2"


def test_pptx_complex_content_judge_scores_diagram_descriptions(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "complex.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                  <p:cSld><p:spTree>
                    <p:sp><p:nvSpPr><p:cNvPr id="1" name="Architecture diagram"
                      descr="Four services connect through the API gateway before storage" /></p:nvSpPr></p:sp>
                    <p:sp><p:nvSpPr><p:cNvPr id="2" name="Risk graph"
                      descr="Graph" /></p:nvSpPr></p:sp>
                  </p:spTree></p:cSld>
                </p:sld>
            """,
        },
    )

    score = PPTXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.sample_findings[0]["name"] == "Risk graph"
    assert score.sample_findings[0]["description"] == "Graph"


def test_pptx_complex_content_judge_scores_equation_context(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "equations.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld
                    xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                    xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math">
                  <p:cSld><p:spTree>
                    <p:sp><p:txBody>
                      <a:p>
                        <a:r><a:t>Forecast model equation for quarterly revenue</a:t></a:r>
                        <m:oMath><m:r><m:t>r = q / t</m:t></m:r></m:oMath>
                      </a:p>
                    </p:txBody></p:sp>
                    <p:sp><p:txBody>
                      <a:p>
                        <m:oMath><m:r><m:t>x = y</m:t></m:r></m:oMath>
                      </a:p>
                    </p:txBody></p:sp>
                  </p:spTree></p:cSld>
                </p:sld>
            """,
        },
    )

    score = PPTXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.per_criterion["equation_context"] == 0.5
    assert score.sample_findings[0]["issue"] == "thin_complex_content_description"
    assert score.sample_findings[0]["kind"] == "equation"
    assert score.sample_findings[0]["source"] == "ppt/slides/slide1.xml#equation-2"


def test_xlsx_complex_content_judge_scores_chart_descriptions(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "complex.xlsx",
        {
            "xl/drawings/drawing1.xml": """
                <xdr:wsDr
                    xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="1" name="Margin chart" descr="Margin improved from 12% to 17%" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                  <xdr:twoCellAnchor>
                    <xdr:pic><xdr:nvPicPr>
                      <xdr:cNvPr id="2" name="Forecast graph" descr="Graph" />
                    </xdr:nvPicPr></xdr:pic>
                  </xdr:twoCellAnchor>
                </xdr:wsDr>
            """,
        },
    )

    score = XLSXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.sample_findings[0]["name"] == "Forecast graph"
    assert score.sample_findings[0]["description"] == "Graph"


def test_xlsx_complex_content_judge_scores_formula_context(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Model"
    worksheet["A1"] = "Revenue"
    worksheet["B1"] = "Total"
    worksheet["A2"] = "Q4"
    worksheet["B2"] = "=SUM(B3:B4)"
    worksheet["D4"] = "=SUM(D1:D3)"
    fixture = tmp_path / "formulas.xlsx"
    workbook.save(fixture)

    score = XLSXComplexContentJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.per_criterion["formula_context"] == 0.5
    assert score.sample_findings[0]["issue"] == "thin_complex_content_description"
    assert score.sample_findings[0]["kind"] == "formula"
    assert score.sample_findings[0]["name"] == "Formula Model!D4"


def test_docx_link_text_judge_reads_hyperlink_relationships(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "links.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                  <w:body>
                    <w:p><w:hyperlink r:id="rId1"><w:r><w:t>click here</w:t></w:r></w:hyperlink></w:p>
                    <w:p><w:hyperlink r:id="rId2"><w:r><w:t>Download annual report</w:t></w:r></w:hyperlink></w:p>
                  </w:body>
                </w:document>
            """,
            "word/_rels/document.xml.rels": """
                <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
                  <Relationship Id="rId1" Target="https://example.com/a" />
                  <Relationship Id="rId2" Target="https://example.com/report" />
                </Relationships>
            """,
        },
    )

    score = DOCXLinkTextJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.60
    assert score.sample_findings[0]["text"] == "click here"
    assert score.sample_findings[0]["target"] == "https://example.com/a"


def test_pptx_link_text_judge_reads_slide_hyperlinks(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "links.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                  <p:cSld><p:spTree><p:sp><p:txBody>
                    <a:p><a:r><a:rPr><a:hlinkClick r:id="rId1" /></a:rPr><a:t>Read more</a:t></a:r></a:p>
                    <a:p><a:r><a:rPr><a:hlinkClick r:id="rId2" /></a:rPr><a:t>Q4 forecast details</a:t></a:r></a:p>
                  </p:txBody></p:sp></p:spTree></p:cSld>
                </p:sld>
            """,
            "ppt/slides/_rels/slide1.xml.rels": """
                <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
                  <Relationship Id="rId1" Target="https://example.com/more" />
                  <Relationship Id="rId2" Target="https://example.com/q4" />
                </Relationships>
            """,
        },
    )

    score = PPTXLinkTextJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.sample_findings[0]["text"] == "Read more"
    assert score.sample_findings[0]["source"] == "ppt/slides/slide1.xml"


def test_pptx_table_structure_judge_reads_table_shapes(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Q2"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$1M"
    fixture = tmp_path / "table.pptx"
    presentation.save(fixture)

    score = PPTXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.confidence == 0.65
    assert score.per_criterion["pptx_table_header_row_presence"] == 1.0
    assert score.per_criterion["table_header_cells_present"] == 0.0
    assert score.sample_findings[0]["issue"] == "pptx_table_empty_header_cells"
    assert score.sample_findings[0]["empty_header_columns"] == [1]


def test_pptx_table_structure_judge_detects_missing_header_row_semantics(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.first_row = False
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "$1M"
    fixture = tmp_path / "missing-header-row.pptx"
    presentation.save(fixture)

    score = PPTXTableStructureJudge(_config()).judge(fixture)

    assert score.score == 0.0
    assert score.per_criterion["pptx_table_header_row_presence"] == 0.0
    assert score.per_criterion["table_header_cells_present"] == 1.0
    assert score.sample_findings[0]["issue"] == "pptx_table_missing_header_row"
    assert score.sample_findings[0]["has_header_row"] is False


def test_xlsx_link_text_judge_reads_worksheet_hyperlinks(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "links.xlsx",
        {
            "xl/worksheets/sheet1.xml": """
                <worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                  <hyperlinks>
                    <hyperlink ref="A1" r:id="rId1" display="https://example.com/raw" />
                    <hyperlink ref="A2" r:id="rId2" display="Quarterly model assumptions" />
                  </hyperlinks>
                </worksheet>
            """,
            "xl/worksheets/_rels/sheet1.xml.rels": """
                <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
                  <Relationship Id="rId1" Target="https://example.com/raw" />
                  <Relationship Id="rId2" Target="https://example.com/model" />
                </Relationships>
            """,
        },
    )

    score = XLSXLinkTextJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.sample_findings[0]["text"] == "https://example.com/raw"
    assert score.sample_findings[0]["target"] == "https://example.com/raw"


def test_docx_decorative_judge_reads_decorative_flags(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative.docx",
        {
            "word/document.xml": """
                <w:document
                    xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                    xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
                    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                    xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
                  <wp:docPr id="1" name="Border">
                    <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
                      <adec:decorative val="1" />
                    </a:ext></a:extLst>
                  </wp:docPr>
                  <wp:docPr id="2" name="Logo" descr="Company logo">
                    <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
                      <adec:decorative val="1" />
                    </a:ext></a:extLst>
                  </wp:docPr>
                </w:document>
            """,
        },
    )

    score = DOCXDecorativeJudge(_config()).judge(fixture)

    assert score.score == 0.5
    assert score.confidence == 0.55
    assert score.sample_findings[0]["issue"] == "decorative_shape_has_accessible_text"
    assert score.sample_findings[0]["description"] == "Company logo"


def test_pptx_decorative_judge_reads_decorative_flags(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative.pptx",
        {
            "ppt/slides/slide1.xml": """
                <p:sld
                    xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                    xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
                  <p:cSld><p:spTree>
                    <p:sp><p:nvSpPr><p:cNvPr id="1" name="Divider">
                      <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
                        <adec:decorative val="1" />
                      </a:ext></a:extLst>
                    </p:cNvPr></p:nvSpPr></p:sp>
                  </p:spTree></p:cSld>
                </p:sld>
            """,
        },
    )

    score = PPTXDecorativeJudge(_config()).judge(fixture)

    assert score.score == 1.0
    assert score.confidence == 0.55
    assert score.sample_findings == []


def test_office_audit_delegates_to_format_specific_judges(monkeypatch) -> None:
    def fake_run_office_checker(file_path: Path, file_type: FileType) -> OfficeCheckReport:  # noqa: ARG001
        return _report(
            file_type,
            [
                OfficeCheckResult("docx-alt-text", "Images contain alternate text", "Passed"),
                OfficeCheckResult("docx-headings", "Headings exist", "Passed"),
                OfficeCheckResult("docx-table-headers", "Tables have headers", "Passed"),
            ],
        )

    def fake_run_office_screen_reader_checks(
        file_path: Path,
        file_type: FileType,
    ) -> OfficeScreenReaderResult:
        return OfficeScreenReaderResult(file_path=file_path, file_type=file_type, issues=[])

    monkeypatch.setattr(office_audit, "run_office_checker", fake_run_office_checker)
    monkeypatch.setattr(
        office_audit,
        "run_office_screen_reader_checks",
        fake_run_office_screen_reader_checks,
    )

    result = office_audit.audit_office_quality(
        Path("fixture.docx"),
        file_type=FileType.DOCX,
        config=_pipeline_config(),
    )

    assert result.format == "docx"
    assert tuple(result.dimensions) == DIMENSIONS_BY_FORMAT["docx"]
    assert sorted(result.behavioral) == [
        "alt_text_substitution",
        "decorative_skip",
        "heading_navigation",
        "reading_order_comprehension",
        "screen_reader_transcript_analysis",
        "table_cell_lookup",
    ]
    assert result.behavioral["alt_text_substitution"].dimension == "alt_text"
    assert (
        result.behavioral["alt_text_substitution"].metadata["behavioral_model"]
        == "qwen2.5:7b"
    )
    assert result.overall_pass is True
    assert "reading_order" not in result.not_applicable_dimensions
    assert "slide_title" in result.not_applicable_dimensions


def test_office_judge_instantiation_enforces_model_separation() -> None:
    bad_config = QualityJudgeConfig(
        backend="ollama",
        model="gemma4:9b",
        production_models=("gemma4:31b-cloud",),
    )

    with pytest.raises(ModelSeparationError):
        DOCXAltTextQualityJudge(bad_config)


def test_office_judge_pairwise_compare_uses_separate_checker_reports() -> None:
    good = _report(
        FileType.DOCX,
        [OfficeCheckResult("docx-alt-text", "Images contain alternate text", "Passed")],
    )
    bad = _report(
        FileType.DOCX,
        [OfficeCheckResult("docx-alt-text", "Images contain alternate text", "Failed")],
    )

    result = DOCXAltTextQualityJudge(_config()).compare(
        Path("a.docx"),
        Path("b.docx"),
        checker_report_a=good,
        checker_report_b=bad,
    )

    assert result == "A_better"


def test_office_judge_prompts_are_version_controlled_files() -> None:
    prompt_expectations = {
        Path("src/project_remedy/quality_judges/office/docx/prompts"): [
            DOCXAltTextQualityJudge,
            DOCXReadingOrderJudge,
            DOCXHeadingSemanticsJudge,
            DOCXTableStructureJudge,
            DOCXLinkTextJudge,
            DOCXDecorativeJudge,
            DOCXComplexContentJudge,
        ],
        Path("src/project_remedy/quality_judges/office/pptx/prompts"): [
            PPTXAltTextQualityJudge,
            PPTXSlideReadingOrderJudge,
            PPTXHeadingSemanticsJudge,
            PPTXTableStructureJudge,
            PPTXLinkTextJudge,
            PPTXDecorativeJudge,
            PPTXComplexContentJudge,
            PPTXSlideTitleJudge,
        ],
        Path("src/project_remedy/quality_judges/office/xlsx/prompts"): [
            XLSXAltTextQualityJudge,
            XLSXTableStructureJudge,
            XLSXLinkTextJudge,
            XLSXComplexContentJudge,
            XLSXSheetOrganizationJudge,
        ],
    }

    for prompt_dir, judges in prompt_expectations.items():
        for judge_cls in judges:
            assert (prompt_dir / judge_cls.prompt_name).exists()
