from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

import pytest

from project_remedy.behavioral_proxies.office.docx import (
    DOCXAltTextSubstitutionTest,
    DOCXDecorativeSkipTest,
    DOCXHeadingNavigationTest,
    DOCXReadingOrderComprehensionTest,
    DOCXScreenReaderTranscriptAnalyzer,
    DOCXTableCellLookupTest,
)
from project_remedy.behavioral_proxies.office.pptx import (
    PPTXAltTextSubstitutionTest,
    PPTXDecorativeSkipTest,
    PPTXScreenReaderTranscriptAnalyzer,
    PPTXSlideReadingOrderComprehensionTest,
    PPTXSlideTitleNavigationTest,
    PPTXTableCellLookupTest,
)
from project_remedy.behavioral_proxies.office.xlsx import (
    XLSXAltTextSubstitutionTest,
    XLSXSheetNavigationTest,
    XLSXTableCellLookupTest,
)
from project_remedy.models import FileType
from project_remedy.office_acceptance import OfficeCheckReport, OfficeCheckResult
from project_remedy.office_acceptance import OfficeScreenReaderIssue, OfficeScreenReaderResult


class _EchoContextAnswerer:
    def answer(self, *, question: str, context: str) -> str:  # noqa: ARG002
        return context


class _RevenueOnlyAnswerer:
    def answer(self, *, question: str, context: str) -> str:  # noqa: ARG002
        if "revenue growth" in context:
            return context
        return "I cannot determine the answer."


def _report(file_type: FileType, results: list[OfficeCheckResult]) -> OfficeCheckReport:
    return OfficeCheckReport(
        file_path=Path("fixture"),
        file_type=file_type,
        results=results,
    )


def _zip_fixture(tmp_path: Path, filename: str, part_name: str, content: str) -> Path:
    path = tmp_path / filename
    with ZipFile(path, "w") as package:
        package.writestr(part_name, content)
    return path


def _mark_docx_header_row(table, value: str = "true") -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    tr_pr = table.rows[0]._tr.get_or_add_trPr()
    tbl_header = OxmlElement("w:tblHeader")
    tbl_header.set(qn("w:val"), value)
    tr_pr.append(tbl_header)


def test_docx_behavioral_proxies_use_office_checker_rules() -> None:
    report = _report(
        FileType.DOCX,
        [
            OfficeCheckResult("docx-alt-text", "Images contain alternate text", "Failed"),
            OfficeCheckResult("docx-headings", "Headings exist", "Passed"),
            OfficeCheckResult("docx-table-headers", "Tables have headers", "Passed"),
        ],
    )

    assert DOCXAltTextSubstitutionTest().run(Path("fixture.docx"), checker_report=report).passed is False
    reading = DOCXReadingOrderComprehensionTest().run(Path("fixture.docx"), checker_report=report)
    assert reading.dimension == "reading_order"
    assert reading.metadata["partial"] is True
    assert DOCXHeadingNavigationTest().run(Path("fixture.docx"), checker_report=report).passed is True
    assert DOCXTableCellLookupTest().run(Path("fixture.docx"), checker_report=report).passed is True


def test_docx_reading_order_proxy_can_use_injected_answerer(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_paragraph("The memo reports revenue growth over five years.")
    fixture = tmp_path / "retained-order.docx"
    document.save(fixture)

    result = DOCXReadingOrderComprehensionTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The memo reports revenue growth over five years.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["parser_support"] == "python_docx_linear_text"
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_docx_reading_order_proxy_flags_answer_retention_loss(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_paragraph("Executive Summary.")
    fixture = tmp_path / "missing-content.docx"
    document.save(fixture)

    result = DOCXReadingOrderComprehensionTest().run(
        fixture,
        answerer=_RevenueOnlyAnswerer(),
        baseline_text="The memo reports revenue growth over five years.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert any(
        finding["issue"] == "llm_answer_retention_loss"
        for finding in result.findings
    )


def test_docx_decorative_skip_reads_ooxml_flags(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
          <wp:docPr id="1" name="Flourish">
            <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
              <adec:decorative val="1" />
            </a:ext></a:extLst>
          </wp:docPr>
          <wp:docPr id="2" name="Logo" descr="Company logo">
            <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
              <adec:decorative val="1" />
            </a:ext></a:extLst>
          </wp:docPr>
        </w:document>
        """,
    )

    result = DOCXDecorativeSkipTest().run(fixture)

    assert result.dimension == "decorative"
    assert result.passed is False
    assert result.score == 0.5
    assert result.metadata["decorative_shape_count"] == 2
    assert result.findings[0]["issue"] == "decorative_shape_has_accessible_text"


def test_office_decorative_skip_rejects_invalid_thresholds(tmp_path: Path) -> None:
    docx_fixture = _zip_fixture(
        tmp_path,
        "decorative-threshold.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
        </w:document>
        """,
    )
    pptx_fixture = _zip_fixture(
        tmp_path,
        "decorative-threshold.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld
            xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        </p:sld>
        """,
    )

    for proxy, fixture in (
        (DOCXDecorativeSkipTest(), docx_fixture),
        (PPTXDecorativeSkipTest(), pptx_fixture),
    ):
        for value, expected in (
            (True, "threshold must be numeric"),
            (float("nan"), "threshold must be finite"),
            (-0.1, "threshold must be between 0 and 1"),
        ):
            try:
                proxy.run(fixture, threshold=value)
            except ValueError as exc:
                assert expected in str(exc)
            else:
                raise AssertionError(
                    f"{proxy.test_name} should reject invalid threshold"
                )


def test_docx_decorative_skip_can_use_injected_answerer(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative-retained.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
          <wp:docPr id="1" name="Divider">
            <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
              <adec:decorative val="1" />
            </a:ext></a:extLst>
          </wp:docPr>
        </w:document>
        """,
    )
    context = "The decorative divider uses teal diagonal lines in the header."

    result = DOCXDecorativeSkipTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text=context,
        candidate_text=context,
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_docx_alt_text_substitution_reads_ooxml_drawing_metadata(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
          <wp:docPr id="1" name="Picture 1" descr="Customer journey map" />
          <wp:docPr id="2" name="Picture 2" />
        </w:document>
        """,
    )

    result = DOCXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.5
    assert result.metadata["parser_support"] == "docx_ooxml_drawing_alt_text"
    assert result.findings[0]["issue"] == "office_object_missing_alt_text"
    assert result.findings[0]["name"] == "Picture 2"


def test_docx_alt_text_substitution_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-generic.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
          <wp:docPr id="1" name="Picture 1" descr="Image" />
          <wp:docPr id="2" name="Picture 2"
            descr="Customer journey map shows 42% renewal path" />
          <wp:docPr id="3" name="Picture 3"
            descr="Customer journey map shows 42% renewal path" />
        </w:document>
        """,
    )

    result = DOCXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.3333
    assert result.metadata["missing_alt_text_count"] == 0
    assert result.metadata["non_substitutive_alt_text_count"] == 1
    assert result.metadata["duplicate_alt_text_count"] == 1
    assert [finding["issue"] for finding in result.findings] == [
        "office_object_non_substitutive_alt_text",
        "office_object_duplicated_substitutive_alt_text",
    ]


def test_docx_alt_text_substitution_can_use_injected_answerer(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-retained.docx",
        "word/document.xml",
        """
        <w:document
            xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
            xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing">
          <wp:docPr id="1" name="Picture 1"
            descr="The customer journey map shows renewal steps." />
        </w:document>
        """,
    )

    result = DOCXAltTextSubstitutionTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The customer journey map shows renewal steps.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_docx_heading_navigation_reads_word_heading_styles(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Overview", level=1)
    document.add_heading("Skipped Details", level=3)
    fixture = tmp_path / "headings.docx"
    document.save(fixture)

    result = DOCXHeadingNavigationTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.5
    assert result.metadata["parser_support"] == "python_docx_heading_styles"
    assert result.findings[0]["issue"] == "docx_heading_level_skip"
    assert result.findings[0]["previous_level"] == 1


def test_docx_heading_navigation_flags_visual_heading_without_semantic_style(
    tmp_path: Path,
) -> None:
    docx = pytest.importorskip("docx")
    from docx.shared import Pt

    document = docx.Document()
    paragraph = document.add_paragraph()
    run = paragraph.add_run("Financial Overview")
    run.bold = True
    run.font.size = Pt(18)
    document.add_paragraph("The financial trend improved over five years.")
    fixture = tmp_path / "visual-heading.docx"
    document.save(fixture)

    result = DOCXHeadingNavigationTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["visual_heading_candidate_count"] == 1
    assert result.findings[0]["issue"] == "docx_visual_heading_without_semantic_style"
    assert result.findings[0]["reasons"] == ["large_text", "bold_text"]


def test_docx_heading_navigation_can_use_injected_answerer(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Enrollment", level=1)
    document.add_paragraph("The enrollment trend improved over five years.")
    fixture = tmp_path / "heading-retained.docx"
    document.save(fixture)

    result = DOCXHeadingNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["navigation_question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_docx_heading_navigation_flags_answer_retention_loss(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Enrollment", level=1)
    document.add_paragraph("The enrollment trend improved over five years.")
    fixture = tmp_path / "heading-loss.docx"
    document.save(fixture)

    result = DOCXHeadingNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="Enrollment",
        candidate_text="candidate missing heading",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_docx_table_cell_lookup_reads_python_docx_tables(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "$1M"
    _mark_docx_header_row(table)
    fixture = tmp_path / "table.docx"
    document.save(fixture)

    result = DOCXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.metadata["parser_support"] == "python_docx_tables"
    assert result.metadata["table_count"] == 1
    assert result.findings[0]["issue"] == "docx_table_empty_header_cells"
    assert result.findings[0]["empty_header_columns"] == [1]


def test_docx_table_cell_lookup_flags_missing_repeated_header_row(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "Region"
    table.cell(1, 0).text = "West"
    fixture = tmp_path / "missing-header.docx"
    document.save(fixture)

    result = DOCXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.findings[0]["issue"] == "docx_table_missing_repeated_header_row"


def test_docx_table_cell_lookup_rejects_false_repeated_header_flag(
    tmp_path: Path,
) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "Region"
    table.cell(1, 0).text = "West"
    _mark_docx_header_row(table, value="false")
    fixture = tmp_path / "false-header.docx"
    document.save(fixture)

    result = DOCXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.findings[0]["issue"] == "docx_table_missing_repeated_header_row"
    assert result.findings[0]["has_header_row"] is False


def test_docx_table_cell_lookup_can_use_injected_answerer(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "120"
    _mark_docx_header_row(table)
    fixture = tmp_path / "lookup-retained.docx"
    document.save(fixture)

    result = DOCXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["lookup_question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_docx_table_cell_lookup_flags_answer_retention_loss(tmp_path: Path) -> None:
    docx = pytest.importorskip("docx")

    document = docx.Document()
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "120"
    _mark_docx_header_row(table)
    fixture = tmp_path / "lookup-loss.docx"
    document.save(fixture)

    result = DOCXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
        candidate_text="candidate missing table value.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_office_screen_reader_transcript_analysis_is_advisory() -> None:
    screen_reader_result = OfficeScreenReaderResult(
        file_path=Path("fixture.docx"),
        file_type=FileType.DOCX,
        issues=[
            OfficeScreenReaderIssue(
                rule_id="docx-alt-text",
                severity="error",
                element="document",
                description="Images contain alternate text",
            )
        ],
    )

    result = DOCXScreenReaderTranscriptAnalyzer().run(
        Path("fixture.docx"),
        screen_reader_result=screen_reader_result,
    )

    assert result.test_name == "screen_reader_transcript_analysis"
    assert result.dimension == "reading_order"
    assert result.passed is False
    assert result.metadata["advisory_only"] is True
    assert result.metadata["error_count"] == 1
    assert result.findings[0]["rule_id"] == "docx-alt-text"


def test_office_screen_reader_transcript_analysis_accepts_provided_transcript() -> None:
    screen_reader_result = OfficeScreenReaderResult(
        file_path=Path("fixture.pptx"),
        file_type=FileType.PPTX,
        issues=[],
    )

    result = PPTXScreenReaderTranscriptAnalyzer().run(
        Path("fixture.pptx"),
        screen_reader_result=screen_reader_result,
        transcript_text="Picture",
    )

    assert result.passed is False
    assert result.metadata["transcript_sources"] == [
        "office_acceptance_screen_reader_checks",
        "pptx_provided_screen_reader_transcript",
    ]
    assert result.findings == [
        {
            "severity": "error",
            "issue": "unlabeled_object_announcement",
            "message": "Transcript announces an object without accessible text.",
            "line_index": 1,
            "announcement": "Picture",
            "source": "pptx_provided_screen_reader_transcript",
        }
    ]


def test_pptx_behavioral_proxies_include_slide_title_and_reading_order_scaffold() -> None:
    report = _report(
        FileType.PPTX,
        [
            OfficeCheckResult("pptx-alt-text", "Pictures contain alternate text", "Passed"),
            OfficeCheckResult("pptx-slide-titles", "Slides expose titles", "Passed"),
        ],
    )

    assert PPTXAltTextSubstitutionTest().run(Path("fixture.pptx"), checker_report=report).passed is True
    assert PPTXSlideTitleNavigationTest().run(Path("fixture.pptx"), checker_report=report).dimension == "slide_title"
    reading = PPTXSlideReadingOrderComprehensionTest().run(Path("fixture.pptx"), slide_count=2)
    assert reading.dimension == "reading_order"
    assert reading.metadata["llm_answering_enabled"] is False
    assert reading.metadata["per_slide"] == [
        {
            "slide_index": 1,
            "passed": True,
            "score": 1.0,
            "llm_answering_enabled": False,
            "parser_support": "python_pptx_shape_order",
            "issue": "parser_unavailable",
            "title_text": "",
            "first_object_text": "",
            "previous_object_text": "",
            "out_of_order_object_text": "",
            "object_count": 0,
            "serialized_text": "",
            "shape_order_texts": [],
            "visual_order_texts": [],
        },
        {
            "slide_index": 2,
            "passed": True,
            "score": 1.0,
            "llm_answering_enabled": False,
            "parser_support": "python_pptx_shape_order",
            "issue": "parser_unavailable",
            "title_text": "",
            "first_object_text": "",
            "previous_object_text": "",
            "out_of_order_object_text": "",
            "object_count": 0,
            "serialized_text": "",
            "shape_order_texts": [],
            "visual_order_texts": [],
        },
    ]


def test_pptx_per_slide_behavioral_proxies_reject_invalid_slide_count_metadata() -> None:
    proxies = [
        PPTXSlideReadingOrderComprehensionTest(),
        PPTXSlideTitleNavigationTest(),
    ]

    for proxy in proxies:
        for value in (True, "2", 1.5, -1):
            with pytest.raises(
                ValueError,
                match="slide_count must be a non-negative integer",
            ):
                proxy.run(Path("fixture.pptx"), slide_count=value)


def test_pptx_reading_order_proxy_flags_title_after_body_shape(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Quarterly Results"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "Body first"
    title_element = title._element
    title_element.getparent().remove(title_element)
    slide.shapes._spTree.append(title_element)
    fixture = tmp_path / "bad-order.pptx"
    presentation.save(fixture)

    result = PPTXSlideReadingOrderComprehensionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.0
    assert result.findings[0]["issue"] == "slide_title_not_first_in_shape_order"
    assert result.metadata["per_slide"][0]["first_object_text"] == "Body first"
    assert result.metadata["per_slide"][0]["object_count"] == 2


def test_pptx_reading_order_proxy_flags_visual_backtracking_in_shape_order(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Results"
    lower = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.5))
    lower.text_frame.text = "Q4 margin details"
    upper = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
    upper.text_frame.text = "Q3 revenue details"
    fixture = tmp_path / "visual-backtracking.pptx"
    presentation.save(fixture)

    result = PPTXSlideReadingOrderComprehensionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.5
    assert result.findings[0]["issue"] == "slide_shape_order_visual_backtracking"
    assert result.findings[0]["previous_object_text"] == "Q4 margin details"
    assert result.findings[0]["out_of_order_object_text"] == "Q3 revenue details"
    assert result.metadata["per_slide"][0]["shape_order_texts"] == [
        "Quarterly Results",
        "Q4 margin details",
        "Q3 revenue details",
    ]
    assert result.metadata["per_slide"][0]["visual_order_texts"] == [
        "Quarterly Results",
        "Q3 revenue details",
        "Q4 margin details",
    ]


def test_pptx_reading_order_proxy_can_use_injected_answerer(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Results"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "The slide reports revenue growth over five years."
    fixture = tmp_path / "retained-order.pptx"
    presentation.save(fixture)

    result = PPTXSlideReadingOrderComprehensionTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_per_slide_text={
            1: "The slide reports revenue growth over five years."
        },
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["per_slide"][0]["answer_accuracy_retention"] == 1.0
    assert result.metadata["per_slide"][0]["candidate_accuracy"] == 1.0
    assert "Quarterly Results" in result.metadata["per_slide"][0]["serialized_text"]


def test_pptx_reading_order_proxy_flags_answer_retention_loss(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Results"
    fixture = tmp_path / "missing-content.pptx"
    presentation.save(fixture)

    result = PPTXSlideReadingOrderComprehensionTest().run(
        fixture,
        answerer=_RevenueOnlyAnswerer(),
        baseline_per_slide_text={
            1: "The slide reports revenue growth over five years."
        },
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["per_slide"][0]["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_pptx_reading_order_proxy_serializes_table_shape_text(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Results"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "120"
    fixture = tmp_path / "table-order.pptx"
    presentation.save(fixture)

    result = PPTXSlideReadingOrderComprehensionTest().run(fixture)

    per_slide = result.metadata["per_slide"][0]
    assert per_slide["object_count"] == 2
    assert "table: Region | Revenue / West | 120" in per_slide["serialized_text"]


def test_pptx_alt_text_substitution_reads_ooxml_drawing_metadata(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cSld><p:spTree>
            <p:pic><p:nvPicPr><p:cNvPr id="1" name="Picture 1"
              descr="Market share by segment" /></p:nvPicPr></p:pic>
            <p:pic><p:nvPicPr><p:cNvPr id="2" name="Picture 2" /></p:nvPicPr></p:pic>
          </p:spTree></p:cSld>
        </p:sld>
        """,
    )

    result = PPTXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.5
    assert result.metadata["parser_support"] == "pptx_ooxml_drawing_alt_text"
    assert result.findings[0]["name"] == "Picture 2"


def test_pptx_alt_text_substitution_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-generic.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cSld><p:spTree>
            <p:pic><p:nvPicPr><p:cNvPr id="1" name="Chart 1"
              descr="Chart" /></p:nvPicPr></p:pic>
            <p:pic><p:nvPicPr><p:cNvPr id="2" name="Chart 2"
              descr="Market share rose from 18% to 24%" /></p:nvPicPr></p:pic>
            <p:pic><p:nvPicPr><p:cNvPr id="3" name="Chart 3"
              descr="Market share rose from 18% to 24%" /></p:nvPicPr></p:pic>
          </p:spTree></p:cSld>
        </p:sld>
        """,
    )

    result = PPTXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.3333
    assert result.metadata["non_substitutive_alt_text_count"] == 1
    assert result.metadata["duplicate_alt_text_count"] == 1
    assert [finding["issue"] for finding in result.findings] == [
        "office_object_non_substitutive_alt_text",
        "office_object_duplicated_substitutive_alt_text",
    ]


def test_pptx_alt_text_substitution_flags_answer_retention_loss(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "alt-loss.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cSld><p:spTree>
            <p:pic><p:nvPicPr><p:cNvPr id="1" name="Picture 1"
              descr="Market share by segment" /></p:nvPicPr></p:pic>
          </p:spTree></p:cSld>
        </p:sld>
        """,
    )

    result = PPTXAltTextSubstitutionTest().run(
        fixture,
        answerer=_RevenueOnlyAnswerer(),
        baseline_text="The slide reports revenue growth over five years.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_pptx_slide_title_navigation_flags_duplicate_titles(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")

    presentation = pptx.Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "Quarterly Results"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "Quarterly Results"
    fixture = tmp_path / "duplicate-titles.pptx"
    presentation.save(fixture)

    result = PPTXSlideTitleNavigationTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.6
    assert result.metadata["parser_support"] == "python_pptx_slide_titles"
    assert [finding["issue"] for finding in result.findings] == [
        "duplicate_slide_title",
        "duplicate_slide_title",
    ]


def test_pptx_slide_title_navigation_can_use_injected_answerer(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Enrollment"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "The enrollment trend improved over five years."
    fixture = tmp_path / "title-retained.pptx"
    presentation.save(fixture)

    result = PPTXSlideTitleNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["navigation_question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_pptx_slide_title_navigation_flags_answer_retention_loss(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Enrollment"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "The enrollment trend improved over five years."
    fixture = tmp_path / "title-loss.pptx"
    presentation.save(fixture)

    result = PPTXSlideTitleNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="Enrollment",
        candidate_text="candidate missing slide title",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_pptx_table_cell_lookup_detects_empty_table_headers(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = ""
    table.cell(0, 1).text = "Q2"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$1M"
    fixture = tmp_path / "table.pptx"
    presentation.save(fixture)

    result = PPTXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.metadata["applicable"] is True
    assert result.metadata["table_count"] == 1
    assert result.findings[0]["issue"] == "pptx_table_empty_header_cells"
    assert result.findings[0]["empty_header_columns"] == [1]


def test_pptx_table_cell_lookup_detects_missing_header_row_semantics(
    tmp_path: Path,
) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.first_row = False
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "$1M"
    fixture = tmp_path / "missing-header-row.pptx"
    presentation.save(fixture)

    result = PPTXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.findings[0]["issue"] == "pptx_table_missing_header_row"
    assert result.findings[0]["has_header_row"] is False


def test_pptx_table_cell_lookup_can_use_injected_answerer(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "120"
    fixture = tmp_path / "lookup-retained.pptx"
    presentation.save(fixture)

    result = PPTXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["lookup_question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_pptx_table_cell_lookup_flags_answer_retention_loss(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1),
    ).table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "West"
    table.cell(1, 1).text = "120"
    fixture = tmp_path / "lookup-loss.pptx"
    presentation.save(fixture)

    result = PPTXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
        candidate_text="candidate missing table value.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_pptx_decorative_skip_reads_ooxml_flags(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld
            xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
          <p:cSld><p:spTree>
            <p:sp><p:nvSpPr><p:cNvPr id="1" name="Divider">
              <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
                <adec:decorative val="1" />
              </a:ext></a:extLst>
            </p:cNvPr></p:nvSpPr></p:sp>
          </p:spTree></p:cSld>
        </p:sld>
        """,
    )

    result = PPTXDecorativeSkipTest().run(fixture)

    assert result.passed is True
    assert result.metadata["decorative_shape_count"] == 1
    assert result.findings == []


def test_pptx_decorative_skip_flags_answer_retention_loss(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "decorative-loss.pptx",
        "ppt/slides/slide1.xml",
        """
        <p:sld
            xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
            xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative">
          <p:cSld><p:spTree>
            <p:sp><p:nvSpPr><p:cNvPr id="1" name="Divider">
              <a:extLst><a:ext uri="{C183D7F6-B498-43B3-948B-1728B52AA6E4}">
                <adec:decorative val="1" />
              </a:ext></a:extLst>
            </p:cNvPr></p:nvSpPr></p:sp>
          </p:spTree></p:cSld>
        </p:sld>
        """,
    )

    result = PPTXDecorativeSkipTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The decorative divider uses teal diagonal lines in the header.",
        candidate_text="candidate missing decorative detail.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_xlsx_behavioral_proxies_cover_table_and_sheet_dimensions() -> None:
    report = _report(
        FileType.XLSX,
        [
            OfficeCheckResult("xlsx-header-behaviors", "Header navigation aids", "Failed"),
        ],
    )

    assert XLSXTableCellLookupTest().run(Path("fixture.xlsx"), checker_report=report).passed is False
    sheet_result = XLSXSheetNavigationTest().run(
        Path("fixture.xlsx"),
        sheet_names=["Sheet1", "Q3 Revenue"],
    )
    assert sheet_result.dimension == "sheet_organization"
    assert sheet_result.passed is False
    assert sheet_result.findings[0]["sheet_name"] == "Sheet1"
    alt = XLSXAltTextSubstitutionTest().run(Path("fixture.xlsx"))
    assert alt.dimension == "alt_text"
    assert alt.metadata["applicable"] is False
    assert alt.metadata["parser_support"] == "xlsx_ooxml_drawing_cnvpr"


def test_xlsx_sheet_navigation_rejects_invalid_sheet_name_metadata() -> None:
    for value in ("Sheet1", [1], [""], True):
        with pytest.raises(
            ValueError,
            match="sheet_names must be a list of non-empty strings",
        ):
            XLSXSheetNavigationTest().run(Path("fixture.xlsx"), sheet_names=value)


def test_xlsx_sheet_navigation_uses_workbook_content_terms(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    archive = workbook.create_sheet("Archive")
    archive.append(["Customer", "Renewal"])
    archive.append(["Northwind", "Yes"])
    default = workbook.create_sheet("Sheet3")
    default.append(["Risk", "Score"])
    default.append(["Security", 7])
    fixture = tmp_path / "sheets.xlsx"
    workbook.save(fixture)

    result = XLSXSheetNavigationTest().run(fixture)

    assert result.passed is False
    assert result.metadata["parser_support"] == "openpyxl"
    assert result.metadata["per_sheet"][0]["sheet_name"] == "Q3 Revenue"
    assert result.metadata["per_sheet"][0]["issue"] == ""
    assert [finding["issue"] for finding in result.findings] == [
        "sheet_name_purpose_unclear",
        "non_descriptive_sheet_name",
    ]


def test_xlsx_sheet_navigation_flags_overview_sheet_after_details(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    summary = workbook.create_sheet("Summary")
    summary.append(["Summary", "Revenue"])
    summary.append(["Total", 120000])
    fixture = tmp_path / "summary-after-details.xlsx"
    workbook.save(fixture)

    result = XLSXSheetNavigationTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.75
    assert result.metadata["per_sheet"][1]["issue"] == "overview_sheet_not_first"
    assert result.findings[0]["sheet_name"] == "Summary"


def test_xlsx_sheet_navigation_preserves_multiple_sheet_issues(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    hidden = workbook.create_sheet("Sheet3")
    hidden.append(["Risk", "Score"])
    hidden.append(["Security", 7])
    hidden.sheet_state = "hidden"
    fixture = tmp_path / "hidden-default-sheet.xlsx"
    workbook.save(fixture)

    result = XLSXSheetNavigationTest().run(fixture)

    assert result.passed is False
    assert result.metadata["per_sheet"][1]["issues"] == [
        "non_descriptive_sheet_name",
        "data_sheet_hidden",
    ]
    assert [finding["issue"] for finding in result.findings] == [
        "non_descriptive_sheet_name",
        "data_sheet_hidden",
    ]


def test_xlsx_sheet_navigation_can_use_injected_answerer(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    fixture = tmp_path / "sheet-nav-retained.xlsx"
    workbook.save(fixture)

    result = XLSXSheetNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["navigation_question_count"] == 1
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_xlsx_sheet_navigation_flags_answer_retention_loss(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    revenue = workbook.active
    revenue.title = "Q3 Revenue"
    revenue.append(["Region", "Revenue"])
    revenue.append(["West", 120000])
    fixture = tmp_path / "sheet-nav-loss.xlsx"
    workbook.save(fixture)

    result = XLSXSheetNavigationTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        candidate_text="candidate missing sheet list",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_xlsx_table_cell_lookup_detects_data_ranges_without_excel_tables(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["Region", "Revenue"])
    worksheet.append(["West", 120000])
    fixture = tmp_path / "plain-range.xlsx"
    workbook.save(fixture)

    result = XLSXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.metadata["parser_support"] == "openpyxl_tables"
    assert result.metadata["table_count"] == 0
    assert result.findings[0]["issue"] == "xlsx_data_range_missing_excel_table"
    assert result.findings[0]["sheet_name"] == "Q3 Revenue"


def test_xlsx_table_cell_lookup_detects_empty_excel_table_headers(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.worksheet.table import Table, TableStyleInfo

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["", "Revenue"])
    worksheet.append(["West", 120000])
    table = Table(displayName="RevenueTable", ref="A1:B2")
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True)
    worksheet.add_table(table)
    fixture = tmp_path / "empty-header.xlsx"
    workbook.save(fixture)

    result = XLSXTableCellLookupTest().run(fixture)

    assert result.passed is False
    assert result.metadata["table_count"] == 1
    assert result.findings[0]["issue"] == "xlsx_table_empty_header_cells"
    assert result.findings[0]["empty_header_columns"] == [1]


def test_xlsx_table_cell_lookup_can_use_injected_answerer(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.worksheet.table import Table, TableStyleInfo

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["Region", "Revenue"])
    worksheet.append(["West", 120])
    worksheet.append(["Total", 120])
    table = Table(displayName="RevenueTable", ref="A1:B3")
    table.totalsRowCount = 1
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True)
    worksheet.add_table(table)
    fixture = tmp_path / "lookup-retained.xlsx"
    workbook.save(fixture)

    result = XLSXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["lookup_question_count"] == 2
    assert result.metadata["answer_accuracy_retention"] == 1.0


def test_xlsx_table_cell_lookup_flags_answer_retention_loss(tmp_path: Path) -> None:
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.worksheet.table import Table, TableStyleInfo

    workbook = openpyxl.Workbook()
    worksheet = workbook.active
    worksheet.title = "Q3 Revenue"
    worksheet.append(["Region", "Revenue"])
    worksheet.append(["West", 120])
    worksheet.append(["Total", 120])
    table = Table(displayName="RevenueTable", ref="A1:B3")
    table.totalsRowCount = 1
    table.tableStyleInfo = TableStyleInfo(name="TableStyleMedium9", showRowStripes=True)
    worksheet.add_table(table)
    fixture = tmp_path / "lookup-loss.xlsx"
    workbook.save(fixture)

    result = XLSXTableCellLookupTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The West region revenue is 120.",
        candidate_text="candidate missing table value.",
    )

    assert result.passed is False
    assert result.score == 0.0
    assert result.metadata["answer_accuracy_retention"] == 0.0
    assert result.findings[0]["issue"] == "llm_answer_retention_loss"


def test_xlsx_alt_text_substitution_reads_drawing_objects(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "drawings.xlsx",
        "xl/drawings/drawing1.xml",
        """
        <xdr:wsDr
            xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="1" name="Chart image" descr="Revenue by quarter" />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="2" name="Unlabeled image" />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
        </xdr:wsDr>
        """,
    )

    result = XLSXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.5
    assert result.metadata["drawing_object_count"] == 2
    assert result.findings[0]["issue"] == "xlsx_drawing_missing_alt_text"


def test_xlsx_alt_text_substitution_flags_generic_and_duplicate_descriptions(
    tmp_path: Path,
) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "drawings-generic.xlsx",
        "xl/drawings/drawing1.xml",
        """
        <xdr:wsDr
            xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="1" name="Forecast graph" descr="Graph" />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="2" name="Margin chart"
                descr="Margin improved from 12% to 17%" />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="3" name="Duplicate margin chart"
                descr="Margin improved from 12% to 17%" />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
        </xdr:wsDr>
        """,
    )

    result = XLSXAltTextSubstitutionTest().run(fixture)

    assert result.passed is False
    assert result.score == 0.3333
    assert result.metadata["missing_alt_text_count"] == 0
    assert result.metadata["non_substitutive_alt_text_count"] == 1
    assert result.metadata["duplicate_alt_text_count"] == 1
    assert [finding["issue"] for finding in result.findings] == [
        "xlsx_drawing_non_substitutive_alt_text",
        "xlsx_drawing_duplicated_substitutive_alt_text",
    ]


def test_xlsx_alt_text_substitution_can_use_injected_answerer(tmp_path: Path) -> None:
    fixture = _zip_fixture(
        tmp_path,
        "drawings-retained.xlsx",
        "xl/drawings/drawing1.xml",
        """
        <xdr:wsDr
            xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing">
          <xdr:twoCellAnchor>
            <xdr:pic><xdr:nvPicPr>
              <xdr:cNvPr id="1" name="Chart image"
                descr="The chart shows revenue growth by quarter." />
            </xdr:nvPicPr></xdr:pic>
          </xdr:twoCellAnchor>
        </xdr:wsDr>
        """,
    )

    result = XLSXAltTextSubstitutionTest().run(
        fixture,
        answerer=_EchoContextAnswerer(),
        baseline_text="The chart shows revenue growth by quarter.",
    )

    assert result.passed is True
    assert result.metadata["llm_answering_enabled"] is True
    assert result.metadata["answer_accuracy_retention"] == 1.0
